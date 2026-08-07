"""Shared helpers for extracting text/images out of the end-of-term pptx
decks under assets/8-1, 8-2, 8-3. Used by build_student_pages.py and
build_group_pages.py.
"""

import glob
import html
import io
import os
import re

from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS_DIR = os.path.join(ROOT, 'assets')
VIDEO_DIR = os.path.join(ASSETS_DIR, '機台影片')

GROUPS = [
    {'id': 'a', 'name': 'A組', 'members': ['蔡翔宇', '詹柏軒', '蔡丞皓'], 'leader': '蔡丞皓'},
    {'id': 'b', 'name': 'B組', 'members': ['黃祈翔', '吳婕語', '石展成', '李心恬'], 'leader': '黃祈翔'},
    {'id': 'c', 'name': 'C組', 'members': ['陳予欣', '陳亮希', '楊初淨', '鍾勻瑨'], 'leader': '楊初淨'},
    {'id': 'd', 'name': 'D組', 'members': ['周以潔', '林為樂', '張宥淇', '周昀希'], 'leader': '周以潔'},
    {'id': 'e', 'name': 'E組', 'members': ['曾詠捷', '黃婕霓', '傅蕾棋', '林琍晴'], 'leader': '傅蕾棋'},
    {'id': 'f', 'name': 'F組', 'members': ['周榆庭', '魏右定', '王儀儼', '周于喆'], 'leader': '周榆庭'},
    {'id': 'g', 'name': 'G組', 'members': ['許詠晴', '陳維潔', '呂睿恩', '王念安'], 'leader': '許詠晴'},
    {'id': 'h', 'name': 'H組', 'members': ['蔡寧', '林宇馨', '黃行書'], 'leader': '蔡寧'},
    {'id': 'i', 'name': 'I組', 'members': ['張諄楷', '陳定鴻', '劉億承', '李侑磬'], 'leader': '陳定鴻'},
]

SECTION_INDIVIDUAL = '壹、個人作品'
SECTION_GROUP = '貳、夜市機台介紹'
SECTION_REFLECTION = '參、個人心得反思'
SUBSECTION_ASSIGNMENT = '分工&我主要負責的部分'

# Applied to every extracted paragraph, e.g. to fix a repeated wording choice
# across every student's deck without having to edit 34 source files.
TEXT_SUBSTITUTIONS = {
    '貼上程式碼': '遊戲程式碼',
}


def normalize_text(text):
    return TEXT_SUBSTITUTIONS.get(text, text)


def find_pptx_for(name):
    matches = [
        f for f in glob.glob(os.path.join(ASSETS_DIR, '8-*', '*.pptx'))
        if os.path.splitext(os.path.basename(f))[0].split('_')[-1] == name
    ]
    return matches[0] if matches else None


def find_manual_photo(name):
    """Look for a hand-supplied photo (e.g. assets/0807-補照片/王念安.JPG),
    added later to fill in gaps left by the pptx auto-extraction. Any
    assets/*補照片*/ folder is searched, newest folder wins on conflict.
    Also matches shared photos named like 「李心恬&吳婕語.JPG」."""
    folders = sorted(glob.glob(os.path.join(ASSETS_DIR, '*補照片*')), reverse=True)
    for folder in folders:
        exact = glob.glob(os.path.join(folder, f'{name}.*'))
        if exact:
            return exact[0]
        for f in glob.glob(os.path.join(folder, f'*{name}*.*')):
            tokens = re.split(r'[&,、]', os.path.splitext(os.path.basename(f))[0])
            if name in tokens:
                return f
    return None


def find_manual_group_photos(group_id):
    """Look for hand-supplied group "製作紀錄" photos, named after the group
    letter (e.g. assets/0807-補照片/D.JPG, or C-1.JPG/C-2.JPG for more than
    one). Any assets/*補照片*/ folder is searched; returns a sorted list of
    paths, empty if the group has no manual photos."""
    letter = group_id.upper()
    folders = sorted(glob.glob(os.path.join(ASSETS_DIR, '*補照片*')), reverse=True)
    matches = []
    for folder in folders:
        matches.extend(glob.glob(os.path.join(folder, f'{letter}.*')))
        matches.extend(glob.glob(os.path.join(folder, f'{letter}-*.*')))
    return sorted(set(matches))


def slide_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip().split('\n')[0]
    return ''


def section_bounds(prs):
    """Return (individual_start, group_start, reflection_start) slide indices, 0-based."""
    titles = [slide_title(s) for s in prs.slides]
    start = next(i for i, t in enumerate(titles) if SECTION_INDIVIDUAL in t)
    group_start = next(i for i, t in enumerate(titles) if SECTION_GROUP in t)
    reflection_start = next(i for i, t in enumerate(titles) if SECTION_REFLECTION in t)
    return start, group_start, reflection_start


def shape_sort_key(shape):
    try:
        return (shape.top or 0, shape.left or 0)
    except Exception:
        return (0, 0)


def find_title_shape(slide):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.has_text_frame:
            ph_type = shape.placeholder_format.type
            if ph_type is not None and 'TITLE' in str(ph_type):
                return shape
    for shape in sorted(slide.shapes, key=shape_sort_key):
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape
    return None


def is_code_block(text):
    lines = [l for l in text.strip().split('\n') if l.strip()]
    if len(lines) < 3:
        return False
    signals = sum(
        1 for l in lines
        if re.search(r'[{};]\s*$', l.strip())
        or re.match(r'^\s*(void|int|float|double|boolean|String|import|if|for|while|function|var|let|const|class|public|private)\b', l.strip())
    )
    return signals >= max(3, len(lines) // 2)


def render_text_shapes(shapes, title_text):
    blocks = []
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        paras = []
        for para in shape.text_frame.paragraphs:
            text = normalize_text(''.join(run.text for run in para.runs).strip())
            if text and text != title_text:
                paras.append({'text': text, 'level': para.level or 0})
        if not paras:
            continue

        full_text = '\n'.join(p['text'] for p in paras)
        if is_code_block(full_text):
            blocks.append(f'<pre class="slide-code"><code>{html.escape(full_text)}</code></pre>')
            continue

        parts = []
        i = 0
        while i < len(paras):
            if paras[i]['level'] > 0:
                items = []
                while i < len(paras) and paras[i]['level'] > 0:
                    items.append(f'<li>{html.escape(paras[i]["text"])}</li>')
                    i += 1
                parts.append(f'<ul>{"".join(items)}</ul>')
            else:
                parts.append(f'<p>{html.escape(paras[i]["text"])}</p>')
                i += 1
        blocks.append(''.join(parts))
    return blocks


def extract_slide_content(slide, slide_no, img_dir, img_url_prefix, img_prefix='slide'):
    title_shape = find_title_shape(slide)
    title = title_shape.text_frame.text.strip().split('\n')[0] if title_shape else ''

    other_shapes = [s for s in sorted(slide.shapes, key=shape_sort_key) if s is not title_shape]
    text_blocks = render_text_shapes(other_shapes, title)

    images = []
    n_img = 0
    for shape in other_shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            n_img += 1
            try:
                image = shape.image
            except Exception:
                continue
            fname = f'{img_prefix}-{slide_no}-{n_img}.{image.ext}'
            with open(os.path.join(img_dir, fname), 'wb') as fh:
                fh.write(image.blob)
            images.append(f'{img_url_prefix}/{fname}')

    return {'title': title, 'text_html': ''.join(text_blocks), 'images': images}


def find_group_videos(group_id):
    """Look for gameplay videos in assets/機台影片/, named <LETTER>.<ext> or
    <LETTER>-<n>.<ext> (e.g. A.mov, B-1.mov, B-2.MOV). Returns a list of
    absolute paths, empty if the group hasn't uploaded a video yet."""
    letter = group_id.upper()
    matches = glob.glob(os.path.join(VIDEO_DIR, f'{letter}.*')) + \
        glob.glob(os.path.join(VIDEO_DIR, f'{letter}-*.*'))
    return sorted(set(matches))


EXHIBIT_LABELS = ['機台的簡介', '機台的遊玩方式', '機台的運作原理']
FUTURE_LABEL = '未來展望'
FUTURE_TITLE_MARKERS = ('未來發展', '未來展望')
EXHIBIT_ORDER = EXHIBIT_LABELS + [FUTURE_LABEL]

SIGN_TITLE_MARKER = '招牌設計'
DEVICE_TITLE_MARKER = '互動遊戲裝置'
DEMO_TITLE_MARKER = '示範影片'


def extract_exhibit_text_fields(prs, group_start, reflection_start):
    """Scan a leader's 「貳、夜市機台介紹」 slides for the exhibit-style
    labelled text fields, in exhibition order: 機台的簡介／機台的遊玩方式／
    機台的運作原理／未來展望 (the last one lives on its own 「六、未來發展」
    section-title slide, with no inline label). Text only -- images are
    collected separately into a single gallery via extract_gallery_images()."""
    blocks = []
    for i in range(group_start, reflection_start):
        slide = prs.slides[i]
        paras = []
        for shape in sorted(slide.shapes, key=shape_sort_key):
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in para.runs).strip()
                if text:
                    paras.append(text)
        if not paras:
            continue

        if any(marker in paras[0] for marker in FUTURE_TITLE_MARKERS):
            text = '\n'.join(p for p in paras[1:] if p)
            if text:
                blocks.append({'label': FUTURE_LABEL, 'text': text})
            continue

        label_positions = [
            (idx, label) for idx, text in enumerate(paras)
            for label in EXHIBIT_LABELS if text.startswith(label)
        ]
        for pos, (idx, label) in enumerate(label_positions):
            end = label_positions[pos + 1][0] if pos + 1 < len(label_positions) else len(paras)
            value_paras = list(paras[idx + 1:end])

            m = re.match(rf'{re.escape(label)}[：:]\s*(.+)', paras[idx])
            if m and m.group(1).strip():
                value_paras.insert(0, m.group(1).strip())

            text = '\n'.join(v for v in value_paras if v)
            if text:
                blocks.append({'label': label, 'text': text})

    blocks.sort(key=lambda b: EXHIBIT_ORDER.index(b['label']))
    return blocks


def find_slide_by_title_marker(prs, group_start, reflection_start, marker):
    for i in range(group_start, reflection_start):
        if marker in slide_title(prs.slides[i]):
            return i
    return None


def _blankness_score(blob):
    """Fraction of near-white pixels in an image, downsized for speed.
    High score ~= a mostly-blank code/text screenshot; low score ~= an
    actual photo with real color and shadow variance."""
    try:
        img = Image.open(io.BytesIO(blob)).convert('RGB').resize((48, 48))
    except Exception:
        return 1.0
    pixels = img.getdata()
    light = sum(1 for r, g, b in pixels if r > 235 and g > 235 and b > 235)
    return light / len(pixels)


def pick_device_cover_image(prs, group_start, reflection_start, img_dir, img_url_prefix, exclude_idx=None):
    """Pick the homepage's "機台" cover tile: the most photo-like image
    across the machine-related slides (skipping 分工, 示範影片, and whichever
    slide was already used for the 招牌 cover). Some decks paste an Arduino
    IDE code screenshot into the 互動遊戲裝置 slide instead of a device photo
    -- those score as almost entirely white/blank and lose out to any real
    photo found elsewhere (e.g. the 成果如何進行 process photo)."""
    best_blob, best_ext, best_score = None, None, None
    for i in range(group_start, reflection_start):
        if i == exclude_idx:
            continue
        title = slide_title(prs.slides[i])
        if SUBSECTION_ASSIGNMENT in title or DEMO_TITLE_MARKER in title:
            continue
        for shape in prs.slides[i].shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                image = shape.image
            except Exception:
                continue
            score = _blankness_score(image.blob)
            if best_score is None or score < best_score:
                best_blob, best_ext, best_score = image.blob, image.ext, score

    if best_blob is None:
        return None
    path = os.path.join(img_dir, f'device.{best_ext}')
    with open(path, 'wb') as fh:
        fh.write(best_blob)
    return f'{img_url_prefix}/device.{best_ext}'


def extract_cover_images(prs, group_start, reflection_start, img_dir, img_url_prefix):
    """Pick the group's homepage cover photos: the 招牌設計 slide's photo
    (the sign/nameplate they built) and the most photo-like machine shot
    elsewhere in the section. Either may be missing."""
    cover = {'sign': None, 'device': None}
    sign_idx = find_slide_by_title_marker(prs, group_start, reflection_start, SIGN_TITLE_MARKER)
    if sign_idx is not None:
        path = extract_largest_picture(prs.slides[sign_idx], os.path.join(img_dir, 'sign'))
        if path:
            cover['sign'] = f'{img_url_prefix}/{os.path.basename(path)}'
    cover['device'] = pick_device_cover_image(prs, group_start, reflection_start, img_dir, img_url_prefix, exclude_idx=sign_idx)
    return cover


def extract_gallery_images(prs, group_start, reflection_start, img_dir, img_url_prefix):
    """Collect every machine-related photo in the leader's 「貳、夜市機台介紹」
    slides (skipping the 分工 slide, the old 示範影片 placeholder slide, and
    the 招牌設計 slide -- that one gets its own standalone showcase image
    instead of sitting in this gallery) for a single photo gallery at the
    bottom of the exhibit text."""
    images = []
    for i in range(group_start, reflection_start):
        title = slide_title(prs.slides[i])
        if SUBSECTION_ASSIGNMENT in title or DEMO_TITLE_MARKER in title or SIGN_TITLE_MARKER in title:
            continue
        content = extract_slide_content(prs.slides[i], i - group_start + 1, img_dir, img_url_prefix, img_prefix='gallery')
        images.extend(content['images'])
    return images


def extract_assignment_paragraphs(slide):
    """Return a 分工 slide's body paragraphs (title excluded), in reading
    order, for line-level self/other-member filtering. Title exclusion is
    done by text match rather than shape identity, since python-pptx hands
    back a fresh wrapper object on every `.shapes` access -- an `is`
    comparison against a shape found via a separate call never matches."""
    title_shape = find_title_shape(slide)
    title_text = title_shape.text_frame.text.strip().split('\n')[0] if title_shape else None
    paras = []
    for shape in sorted(slide.shapes, key=shape_sort_key):
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = ''.join(run.text for run in para.runs).strip()
            if text and text != title_text:
                paras.append(text)
    return paras


def find_machine_name(prs, group_start, reflection_start):
    """Scan a leader's 「貳、夜市機台介紹」 slides for the 「機台的名稱：...」
    label and return the quoted name after it, e.g. 「射企球」 -> 射企球."""
    for i in range(group_start, reflection_start):
        paras = []
        for shape in sorted(prs.slides[i].shapes, key=shape_sort_key):
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in para.runs).strip()
                if text:
                    paras.append(text)
        for idx, text in enumerate(paras):
            if '機台的名稱' not in text:
                continue
            m = re.search(r'機台的名稱[：:]\s*(.+)', text)
            candidate = m.group(1).strip() if m and m.group(1).strip() else None
            if not candidate and idx + 1 < len(paras):
                candidate = paras[idx + 1].strip()
            if candidate:
                return candidate.strip('「」『』<>〈〉"\' ')
    return None


def extract_largest_picture(slide, out_path):
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pics:
        return None
    best = max(pics, key=lambda s: (s.width or 0) * (s.height or 0))
    try:
        image = best.image
    except Exception:
        return None
    path = f'{out_path}.{image.ext}'
    with open(path, 'wb') as fh:
        fh.write(image.blob)
    return path


def render_content_heading(title):
    return f'<h2 class="content-heading">{html.escape(title)}</h2>'


def render_slide_blocks(slides, empty_message):
    blocks = []
    for slide in slides:
        if not slide['title'] and not slide['text_html'] and not slide['images']:
            continue

        if slide['title'] and not slide['text_html'] and not slide['images']:
            blocks.append(render_content_heading(slide['title']))
            continue

        if slide.get('heading_style'):
            images_html = ''
            if slide['images']:
                imgs = ''.join(f'<img src="{html.escape(src)}" alt="{html.escape(slide["title"])}" loading="lazy" />' for src in slide['images'])
                images_html = f'<div class="slide-images">{imgs}</div>'
            blocks.append(f'{render_content_heading(slide["title"])}{images_html}')
            continue

        images_html = ''
        if slide['images']:
            imgs = ''.join(f'<img src="{html.escape(src)}" alt="{html.escape(slide["title"])}" loading="lazy" />' for src in slide['images'])
            images_html = f'<div class="slide-images">{imgs}</div>'

        text_html = f'<div class="slide-text">{slide["text_html"]}</div>' if slide['text_html'] else ''
        title_html = f'<h3 class="slide-block__title">{html.escape(slide["title"])}</h3>' if slide['title'] else ''

        blocks.append(f'<section class="slide-block">{title_html}{images_html}{text_html}</section>')

    return '\n    '.join(blocks) if blocks else f'<p class="student-empty">{empty_message}</p>'
