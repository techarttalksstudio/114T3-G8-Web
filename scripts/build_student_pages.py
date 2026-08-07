#!/usr/bin/env python3
"""Extract each student's 「個人作品」section (text + images) from their
end-of-term pptx and render it as a standalone static page under students/.

Re-run any time the source .pptx files under assets/8-1, 8-2, 8-3 change.
Does not touch js/data.js, js/main.js or css/style.css.
"""

import glob
import html
import os
import shutil

from pptx import Presentation

from pptx_common import (
    GROUPS, ASSETS_DIR, find_pptx_for, find_manual_photo, section_bounds,
    extract_slide_content, extract_largest_picture, render_slide_blocks,
)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
MEDIA_DIR = os.path.join(ASSETS_DIR, 'media', 'students')
STUDENTS_DIR = os.path.join(ROOT, 'students')

# Every student's deck repeats the same generic「1.成果如何進行／2.對應能力」
# blurb on these two slides (it's boilerplate from the course template, not
# personal content) — drop the text. The Joystick slide additionally keeps
# only its last image (the wiring-diagram photo), dropping the screenshot.
# Both are section-divider slides, so they render as a big borderless
# heading (optionally with that one image below) instead of a boxed
# slide-block like the rest of the content.
SLIDE_OVERRIDES = {
    '一、Scratch x Joystick': {'strip_text': True, 'keep_last_image_only': True, 'heading_style': True},
    '二、Processing基礎 - 彩色泡泡、射氣球、乒乓球裝置': {'strip_text': True, 'heading_style': True},
}


# Students whose assets/media/students/<key>-photo.<ext> was hand-cropped/
# adjusted after generation — never delete or replace it on rebuild.
PROTECTED_PHOTOS = {'e-1', 'a-1'}


def apply_slide_overrides(slide):
    override = SLIDE_OVERRIDES.get(slide['title'])
    if not override:
        return slide
    if override.get('strip_text'):
        slide['text_html'] = ''
    if override.get('keep_last_image_only') and slide['images']:
        slide['images'] = slide['images'][-1:]
    if override.get('heading_style'):
        slide['heading_style'] = True
    return slide


def build_student_list():
    students = []
    for group in GROUPS:
        for idx, name in enumerate(group['members']):
            students.append({
                'key': f"{group['id']}-{idx + 1}",
                'name': name,
                'group_id': group['id'],
                'group_name': group['name'],
            })
    return students


def render_page(student, slides, photo_url):
    photo_html = (
        f'<img src="{html.escape(photo_url)}" alt="{html.escape(student["name"])}的照片" />'
        if photo_url else
        '<img src="../assets/placeholder/photo-placeholder.svg" alt="尚未提供照片" />'
    )

    blocks_html = render_slide_blocks(slides, '個人作品內容準備中。')

    return f'''<!DOCTYPE html>
<html lang="zh-Hant">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>{html.escape(student["name"])}的個人作品｜台灣夜市互動遊戲成果展</title>
  <link rel="stylesheet" href="../css/style.css" />
</head>
<body>
  <nav class="site-nav">
    <div class="container">
      <a href="../index.html" class="site-nav__brand">夜市互動遊戲成果展</a>
      <ul class="site-nav__links">
        <li><a href="../groups/{student["group_id"]}.html">查看{html.escape(student["group_name"])}夜市機台 →</a></li>
        <li><a href="../index.html#students">← 回個人作品區</a></li>
      </ul>
    </div>
  </nav>

  <header class="student-hero">
    <div class="student-hero__photo">{photo_html}</div>
    <h1>{html.escape(student["name"])}</h1>
    <p class="student-hero__group">{html.escape(student["group_name"])}｜個人作品</p>
  </header>

  <main class="section">
    <div class="container student-content">
    {blocks_html}
    </div>
  </main>

  <footer class="site-footer">
    <p>台灣夜市互動遊戲｜期末成果展示網站</p>
  </footer>

  <script src="../js/lightbox.js"></script>
</body>
</html>
'''


def main():
    os.makedirs(STUDENTS_DIR, exist_ok=True)
    os.makedirs(MEDIA_DIR, exist_ok=True)

    students = build_student_list()
    missing_photo = []
    missing_pptx = []
    missing_sections = []

    for student in students:
        key = student['key']
        pptx_path = find_pptx_for(student['name'])
        if not pptx_path:
            missing_pptx.append(student['name'])
            continue

        prs = Presentation(pptx_path)

        try:
            start, group_start, _ = section_bounds(prs)
        except StopIteration:
            missing_sections.append(student['name'])
            continue

        img_dir = os.path.join(MEDIA_DIR, key)
        os.makedirs(img_dir, exist_ok=True)
        img_url_prefix = f'../assets/media/students/{key}'

        slides = [
            apply_slide_overrides(extract_slide_content(prs.slides[i], i - start + 1, img_dir, img_url_prefix))
            for i in range(start + 1, group_start)  # skip the 「壹、個人作品」section-title slide itself
        ]

        if key in PROTECTED_PHOTOS:
            existing = glob.glob(os.path.join(MEDIA_DIR, f'{key}-photo.*'))
            photo_url = f'../assets/media/students/{os.path.basename(existing[0])}' if existing else None
        else:
            for stale in glob.glob(os.path.join(MEDIA_DIR, f'{key}-photo.*')):
                os.remove(stale)

            photo_url = None
            manual_photo = find_manual_photo(student['name'])
            if manual_photo:
                ext = os.path.splitext(manual_photo)[1].lower()
                dest = os.path.join(MEDIA_DIR, f'{key}-photo{ext}')
                shutil.copyfile(manual_photo, dest)
                photo_url = f'../assets/media/students/{os.path.basename(dest)}'
            else:
                photo_path = extract_largest_picture(prs.slides[len(prs.slides) - 1], os.path.join(MEDIA_DIR, f'{key}-photo'))
                if photo_path:
                    photo_url = f'../assets/media/students/{os.path.basename(photo_path)}'

        if not photo_url:
            missing_photo.append(student['name'])

        html_out = render_page(student, slides, photo_url)
        with open(os.path.join(STUDENTS_DIR, f'{key}.html'), 'w', encoding='utf-8') as fh:
            fh.write(html_out)

        print(f'✓ {key:5s} {student["name"]}  ({len(slides)} slides, photo={"yes" if photo_url else "no"})')

    print()
    print(f'完成 {len(students) - len(missing_pptx) - len(missing_sections)}/{len(students)} 位學生的頁面')
    if missing_pptx:
        print('找不到簡報檔案:', '、'.join(missing_pptx))
    if missing_sections:
        print('找不到章節標題（需人工確認頁碼範圍）:', '、'.join(missing_sections))
    if missing_photo:
        print(f'缺個人照片（{len(missing_photo)} 位，需另外補）:', '、'.join(missing_photo))


if __name__ == '__main__':
    main()
