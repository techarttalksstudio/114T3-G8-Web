#!/usr/bin/env python3
"""Build each group's static page under groups/.

Exhibit-style page: a gameplay video (from assets/機台影片/) up top, then
the machine's 機台的簡介／機台的遊玩方式／機台的運作原理／未來展望 text
pulled out of the *group leader's* deck, followed by a single photo
gallery of every machine-related image. Below that, 小組分工 lists every
member (leader first) with only the line(s) they wrote about their own
part -- other members' lines are filtered out. If a member has no 分工
slide of their own, we fall back to pulling their line out of the
leader's 分工 slide; if that has nothing either, the block is left blank.

Re-run any time the source .pptx files under assets/8-1, 8-2, 8-3 change,
or when assets/機台影片/ gets new videos.
Does not touch js/main.js or css/style.css. js/data.js's per-group
signSrc/deviceSrc cover images should be updated by hand from the paths
this script prints.
"""

import html
import os
import re
import shutil

from pptx import Presentation

from pptx_common import (
    GROUPS, ASSETS_DIR, find_pptx_for, section_bounds, slide_title,
    extract_slide_content, extract_exhibit_text_fields, extract_gallery_images,
    extract_cover_images, extract_assignment_paragraphs, find_group_videos,
    find_manual_group_photos, render_content_heading, find_machine_name,
    SUBSECTION_ASSIGNMENT,
)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
MEDIA_DIR = os.path.join(ASSETS_DIR, 'media', 'groups')
GROUPS_DIR = os.path.join(ROOT, 'groups')

HEADER_WORDS = {'小組分工', '分工', '組員分工', '我負責的部分的圖片'}
NAME_LINE_RE = re.compile(r'^[>🔎●◆\-\*\s]*([^：:]{1,24})[：:]\s*(.+)$')
NAME_TOKEN_RE = re.compile(r'^[一-鿿A-Za-z（）() 、,，&]+$')

# Manual corrections applied after the automatic image extraction, so a
# known-bad auto-picked photo (e.g. one slide's picture duplicating another
# slide's, a code screenshot, or a photo that needed hand-rotating) stays
# fixed across re-runs instead of reverting every time the pptx is
# re-scanned. Maps group id -> {generated filename: replacement source path
# relative to the repo root}.
GALLERY_IMAGE_OVERRIDES = {
    'd': {
        'gallery-3-1.png': 'assets/機台影片/D.png',  # was a duplicate of the sign photo; swapped for the real target-conveyor photo
        'device.png': 'assets/機台影片/D.png',  # same swap for the homepage card's 機台 cover tile
    },
    'c': {
        'sign.png': 'assets/media/groups/c/sign-protected.png',  # hand-rotated to upright; source pptx photo is sideways
    },
}


def apply_gallery_overrides(group_id, img_dir):
    for filename, src_rel in GALLERY_IMAGE_OVERRIDES.get(group_id, {}).items():
        dest = os.path.join(img_dir, filename)
        src = os.path.join(ROOT, src_rel)
        if os.path.exists(dest) and os.path.exists(src):
            shutil.copyfile(src, dest)


def add_manual_gallery_photos(group_id, img_dir, img_url_prefix):
    """Hand-supplied 製作紀錄 photos (assets/*補照片*/<LETTER>.* or
    <LETTER>-<n>.*, e.g. D.JPG or C-1.JPG/C-2.JPG) copied alongside the
    auto-extracted gallery images. Returns their URLs to append."""
    urls = []
    for i, src in enumerate(find_manual_group_photos(group_id), start=1):
        ext = os.path.splitext(src)[1].lower()
        filename = f'extra-{i}{ext}'
        dest = os.path.join(img_dir, filename)
        shutil.copyfile(src, dest)
        urls.append(f'{img_url_prefix}/{filename}')
    return urls


def find_assignment_index(prs, group_start, reflection_start):
    for i in range(group_start, reflection_start):
        if SUBSECTION_ASSIGNMENT in slide_title(prs.slides[i]):
            return i
    return None


def _split_name_tokens(names_part):
    return [t.strip() for t in re.split(r'[、,，&\s]+', names_part) if t.strip()]


def _mentions(tokens, name):
    return any(name in t or t in name for t in tokens)


def filter_own_lines(paras, member_name, keep_freeform=True):
    """Keep only the paragraph(s) a 分工 slide attributes to member_name,
    dropping lines clearly written about a different member. Lines with no
    recognizable "name：desc" shape are freeform prose -- on a member's own
    slide that's assumed to be about themselves and kept, but when this is
    called as a fallback against the *leader's* slide to find some other
    member's line, freeform text carries no name attribution and must be
    dropped (keep_freeform=False), or it would misattribute the leader's
    own prose to whoever else is being looked up."""
    kept = []
    for raw_line in paras:
        line = re.sub(r'^[>🔎●◆\-\*\s]+', '', raw_line.strip())
        if line.strip('： :') in HEADER_WORDS:
            continue

        remainder = line
        while True:
            m = NAME_LINE_RE.match(remainder)
            if not m or m.group(1).strip() not in HEADER_WORDS:
                break
            remainder = m.group(2).strip()

        if not remainder or remainder.strip('： :') in HEADER_WORDS:
            continue

        m = NAME_LINE_RE.match(remainder)
        if not m:
            if keep_freeform:
                kept.append(remainder)
            continue

        names_part, desc = m.group(1).strip(), m.group(2).strip()
        if not NAME_TOKEN_RE.match(names_part):
            if keep_freeform:
                kept.append(remainder)
            continue

        tokens = _split_name_tokens(names_part)
        if _mentions(tokens, member_name):
            if all(member_name in t or t in member_name for t in tokens):
                # Every name in the "Name：" prefix is just this member
                # (possibly a nickname) -- redundant with the block's <h3>
                # heading above it, so drop the prefix and keep only desc.
                kept.append(desc)
            else:
                # Shared line naming other members too (e.g. "吳婕語、李心恬：
                # 招牌") -- keep the full prefix, it's informative.
                kept.append(remainder)
        else:
            # Name-shaped line ("Someone：did X") that isn't self -- drop it
            # even if the name doesn't exactly match another roster member
            # (source decks have typos), since a name-colon-description line
            # is always attributed to *someone*, just not this member.
            continue

    return kept


def render_group_line(group, member_keys):
    member_links = '、'.join(
        f'<a href="../students/{key}.html">{html.escape(name)}</a>'
        for name, key in member_keys
    )
    leader_key = next(key for name, key in member_keys if name == group['leader'])
    leader_link = f'<a href="../students/{leader_key}.html">{html.escape(group["leader"])}</a>'
    return f'<p class="student-hero__group">組員：{member_links}｜組長：{leader_link}</p>'


def render_assignment_blocks(assignment_data):
    items = []
    for entry in assignment_data:
        images_html = ''
        if entry['images']:
            imgs = ''.join(f'<img src="{html.escape(src)}" alt="{html.escape(entry["name"])}" loading="lazy" />' for src in entry['images'])
            images_html = f'<div class="slide-images">{imgs}</div>'
        text_html = ''.join(f'<p>{html.escape(line)}</p>' for line in entry['lines'])
        text_block = f'<div class="slide-text">{text_html}</div>' if text_html else ''
        items.append(
            f'<div class="assignment-block">'
            f'<h3 class="assignment-block__name"><a href="../students/{entry["key"]}.html">{html.escape(entry["name"])}</a></h3>'
            f'{text_block}{images_html}'
            f'</div>'
        )
    return f'<div class="assignment-list">{"".join(items)}</div>'


def render_exhibit_blocks(fields, gallery_images):
    if not fields and not gallery_images:
        return '<p class="student-empty">機台介紹內容準備中。</p>'

    blocks = []
    for field in fields:
        text_html = ''.join(f'<p>{html.escape(line)}</p>' for line in field['text'].split('\n') if line)
        blocks.append(
            f'<section class="exhibit-block">'
            f'<span class="exhibit-block__label">{html.escape(field["label"])}</span>'
            f'<div class="slide-text">{text_html}</div>'
            f'</section>'
        )

    if gallery_images:
        imgs = ''.join(f'<img src="{html.escape(src)}" alt="製作紀錄" loading="lazy" />' for src in gallery_images)
        blocks.append(
            f'<section class="exhibit-block">'
            f'<span class="exhibit-block__label">製作紀錄</span>'
            f'<div class="slide-images slide-images--gallery">{imgs}</div>'
            f'</section>'
        )

    return '\n    '.join(blocks)


IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.webp'}


def render_media_html(media_paths):
    """Render the top-of-page hero + any extras from assets/機台影片/. Most
    groups drop a video there, but a group whose machine isn't finished yet
    can drop a photo instead (e.g. D.png) -- rendered as a plain image hero
    rather than a <video> tag."""
    if not media_paths:
        return '', ''

    def rel(p):
        return f'../assets/機台影片/{os.path.basename(p)}'

    def is_image(p):
        return os.path.splitext(p)[1].lower() in IMAGE_EXTS

    first = media_paths[0]
    if is_image(first):
        hero_html = f'''<div class="student-hero__photo">
      <img src="{html.escape(rel(first))}" alt="機台照片" loading="lazy" />
    </div>'''
    else:
        hero_html = f'''<div class="student-hero__photo student-hero__photo--video">
      <video src="{html.escape(rel(first))}" controls preload="metadata"></video>
    </div>'''

    extra_html = ''
    if len(media_paths) > 1:
        extras = ''.join(
            f'<img src="{html.escape(rel(p))}" alt="機台照片" loading="lazy" class="exhibit-video" />'
            if is_image(p) else
            f'<video src="{html.escape(rel(p))}" controls preload="metadata" class="exhibit-video"></video>'
            for p in media_paths[1:]
        )
        extra_html = f'<section class="exhibit-block"><span class="exhibit-block__label">更多影片</span>{extras}</section>'

    return hero_html, extra_html


def render_sign_showcase(sign_url, group_name):
    if not sign_url:
        return ''
    return (
        '<section class="exhibit-block">'
        '<span class="exhibit-block__label">招牌設計</span>'
        f'<div class="sign-showcase"><img src="{html.escape(sign_url)}" alt="{html.escape(group_name)}招牌" loading="lazy" /></div>'
        '</section>'
    )


def render_page(group, machine_name, video_paths, exhibit_fields, gallery_images, member_keys, assignment_data, sign_url):
    assignment_html = render_assignment_blocks(assignment_data)
    sign_html = render_sign_showcase(sign_url, group['name'])
    exhibit_html = render_exhibit_blocks(exhibit_fields, gallery_images)
    video_hero_html, video_extra_html = render_media_html(video_paths)
    group_line_html = render_group_line(group, member_keys)

    if machine_name:
        heading_html = f'<h1 class="group-hero__title">「{html.escape(machine_name)}」</h1>'
        groupname_html = f'<p class="group-hero__groupname">{html.escape(group["name"])}</p>'
    else:
        heading_html = f'<h1>{html.escape(group["name"])}</h1>'
        groupname_html = ''

    return f'''<!DOCTYPE html>
<html lang="zh-Hant">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>{html.escape(group["name"])}的團體作品｜台灣夜市互動遊戲成果展</title>
  <link rel="stylesheet" href="../css/style.css" />
</head>
<body>
  <nav class="site-nav">
    <div class="container">
      <a href="../index.html" class="site-nav__brand">第一版機台的原型紀錄</a>
      <ul class="site-nav__links">
        <li><a href="../index.html#groups">← 回團體作品區</a></li>
      </ul>
    </div>
  </nav>

  <header class="student-hero">
    {video_hero_html}
    {heading_html}
    {groupname_html}
    {group_line_html}
  </header>

  <main class="section">
    <div class="container student-content">
    {video_extra_html}
    {sign_html}
    {exhibit_html}
    {render_content_heading('小組分工')}
    {assignment_html}
    </div>
  </main>

  <footer class="site-footer">
    <p>台灣夜市互動遊戲｜期末成果展示網站</p>
  </footer>

  <script src="../js/lightbox.js"></script>
</body>
</html>
'''


def build_assignment_data(group, leader_paras, img_dir, img_url_prefix, issues):
    # Leader first, then the rest in roster order.
    ordered_names = [group['leader']] + [m for m in group['members'] if m != group['leader']]

    data = []
    for name in ordered_names:
        key = f'{group["id"]}-{group["members"].index(name) + 1}'

        own_lines = []
        own_images = []
        pptx_path = find_pptx_for(name)
        if not pptx_path:
            issues.append(f'{group["name"]}: 找不到 {name} 的簡報')
        else:
            member_prs = Presentation(pptx_path)
            try:
                _, m_group_start, m_reflection_start = section_bounds(member_prs)
            except StopIteration:
                issues.append(f'{group["name"]}: {name} 的簡報找不到章節標題')
                m_group_start = None

            if m_group_start is not None:
                idx = find_assignment_index(member_prs, m_group_start, m_reflection_start)
                if idx is None:
                    issues.append(f'{group["name"]}: {name} 的簡報找不到分工頁')
                else:
                    own_paras = extract_assignment_paragraphs(member_prs.slides[idx])
                    own_lines = filter_own_lines(own_paras, name)
                    content = extract_slide_content(
                        member_prs.slides[idx], group['members'].index(name) + 1,
                        img_dir, img_url_prefix, img_prefix='assign',
                    )
                    own_images = content['images']

        lines = own_lines
        if not lines and name != group['leader']:
            lines = filter_own_lines(leader_paras, name, keep_freeform=False)

        data.append({'name': name, 'key': key, 'lines': lines, 'images': own_images})

    return data


def main():
    os.makedirs(GROUPS_DIR, exist_ok=True)
    os.makedirs(MEDIA_DIR, exist_ok=True)

    issues = []
    completed = 0

    for group in GROUPS:
        leader_pptx = find_pptx_for(group['leader'])
        if not leader_pptx:
            issues.append(f'{group["name"]}: 找不到組長 {group["leader"]} 的簡報')
            continue

        prs = Presentation(leader_pptx)
        try:
            _, group_start, reflection_start = section_bounds(prs)
        except StopIteration:
            issues.append(f'{group["name"]}: 組長簡報找不到章節標題')
            continue

        img_dir = os.path.join(MEDIA_DIR, group['id'])
        os.makedirs(img_dir, exist_ok=True)
        img_url_prefix = f'../assets/media/groups/{group["id"]}'

        machine_name = find_machine_name(prs, group_start, reflection_start)
        exhibit_fields = extract_exhibit_text_fields(prs, group_start, reflection_start)
        gallery_images = extract_gallery_images(prs, group_start, reflection_start, img_dir, img_url_prefix)
        gallery_images += add_manual_gallery_photos(group['id'], img_dir, img_url_prefix)
        cover = extract_cover_images(prs, group_start, reflection_start, img_dir, img_url_prefix)
        apply_gallery_overrides(group['id'], img_dir)
        video_paths = find_group_videos(group['id'])
        if not video_paths:
            issues.append(f'{group["name"]}: assets/機台影片/ 裡沒有找到影片或圖片')

        leader_assignment_idx = find_assignment_index(prs, group_start, reflection_start)
        leader_paras = extract_assignment_paragraphs(prs.slides[leader_assignment_idx]) if leader_assignment_idx is not None else []

        member_keys = [(name, f'{group["id"]}-{i + 1}') for i, name in enumerate(group['members'])]
        assignment_data = build_assignment_data(group, leader_paras, img_dir, img_url_prefix, issues)

        html_out = render_page(group, machine_name, video_paths, exhibit_fields, gallery_images, member_keys, assignment_data, cover['sign'])
        with open(os.path.join(GROUPS_DIR, f'{group["id"]}.html'), 'w', encoding='utf-8') as fh:
            fh.write(html_out)

        with_lines = sum(1 for e in assignment_data if e['lines'])
        completed += 1
        print(f'✓ {group["id"]}  {group["name"]}（組長 {group["leader"]}）：機台名稱「{machine_name}」，'
              f'{len(video_paths)} 支影片，{len(exhibit_fields)} 個文字欄位，{len(gallery_images)} 張製作紀錄照片，'
              f'{with_lines}/{len(assignment_data)} 位分工，'
              f'招牌照 {cover["sign"]}，機台照 {cover["device"]}')

    print()
    print(f'完成 {completed}/{len(GROUPS)} 組頁面')
    if issues:
        print('需要注意:')
        for i in issues:
            print(' -', i)


if __name__ == '__main__':
    main()
